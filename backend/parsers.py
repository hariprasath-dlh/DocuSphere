import io
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def parse_pdf(contents: bytes) -> list[dict]:
    pdf_reader = PdfReader(io.BytesIO(contents))
    pages_data = []
    for i, page in enumerate(pdf_reader.pages):
        text = page.extract_text()
        if text:
            pages_data.append({"page": i + 1, "text": text})
    return pages_data

def parse_docx(contents: bytes) -> list[dict]:
    doc = Document(io.BytesIO(contents))
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    # Since docx doesn't have "pages" in a standard accessible way easily, 
    # we treat it as one or block-segments.
    return [{"page": 1, "text": "\n".join(full_text)}]

def parse_pptx(contents: bytes) -> list[dict]:
    prs = Presentation(io.BytesIO(contents))
    slides_data = []
    for i, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        if text:
            slides_data.append({"page": i + 1, "text": "\n".join(text)})
    return slides_data

def parse_xlsx(contents: bytes) -> list[dict]:
    df = pd.read_excel(io.BytesIO(contents))
    # Convert whole sheet to text
    text = df.to_string()
    return [{"page": 1, "text": text}]

def parse_txt(contents: bytes) -> list[dict]:
    text = contents.decode("utf-8", errors="ignore")
    return [{"page": 1, "text": text}]

def get_pages_data(contents: bytes, filename: str) -> list[dict]:
    ext = filename.split(".")[-1].lower()
    if ext == "pdf":
        return parse_pdf(contents)
    elif ext == "docx":
        return parse_docx(contents)
    elif ext == "pptx":
        return parse_pptx(contents)
    elif ext in ["xlsx", "xls"]:
        return parse_xlsx(contents)
    elif ext == "txt":
        return parse_txt(contents)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")
